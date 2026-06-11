"""Core logic for the ESG Value Creation Module. Two stages:
   EXTRACT  -> pulls facts from the DD, each with a source; flags gaps.
   ADVISE   -> builds the detailed 100-day roadmap on top of those facts.
"""
import os, json, re
from io import BytesIO
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

TOCONFIRM = "[TO CONFIRM \u2014 not in DD]"

DIMENSIONS = [
    ("risk_reduction", "Risk reduction"), ("ebitda_impact", "EBITDA impact"),
    ("revenue_growth", "Revenue growth"), ("financing_benefit", "Financing benefits"),
    ("operational_efficiency", "Operational efficiency"), ("implementation_feasibility", "Implementation feasibility"),
]
DEFAULT_WEIGHTS = {"risk_reduction": 25, "ebitda_impact": 20, "revenue_growth": 10,
                   "financing_benefit": 10, "operational_efficiency": 15, "implementation_feasibility": 20}
PHASES = [("0-30", "Day 0-30 \u00b7 Stabilise & assign ownership"),
          ("30-60", "Day 30-60 \u00b7 Implement quick wins"),
          ("60-100", "Day 60-100 \u00b7 Embed KPIs & board reporting")]

# Holtara brand
NAVY = RGBColor(0x2C, 0x3C, 0x7E); ORANGE = RGBColor(0xE8, 0x77, 0x22)
GREY = RGBColor(0x57, 0x57, 0x56); TEAL = RGBColor(0x6F, 0xA4, 0x91)
TEXT = RGBColor(0x1E, 0x1E, 0x1E); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
REDFLAG = RGBColor(0xC0, 0x39, 0x2B); LIGHT = RGBColor(0xF4, 0xF4, 0xF2)
FONT = "Open Sans"
PHASE_COLOR = {"0-30": NAVY, "30-60": TEAL, "60-100": ORANGE}
PILLAR_COLOR = {"E": TEAL, "S": NAVY, "G": ORANGE}
LOGO_PATH = "holtara_logo.png"


# ---------------------------------------------------------------- extraction
def extract_pptx_text(file_bytes):
    prs = Presentation(BytesIO(file_bytes)); chunks = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                parts.append("TABLE:\n" + "\n".join(
                    " | ".join(c.text.strip() for c in r.cells) for r in shape.table.rows))
        if parts:
            chunks.append("--- Slide %d ---\n%s" % (i + 1, "\n".join(parts)))
    return "\n\n".join(chunks)


def parse_json(text):
    text = re.sub(r"```$", "", re.sub(r"^```[a-zA-Z]*", "", text.strip()).strip()).strip()
    a, b = text.find("{"), text.rfind("}")
    if a != -1 and b != -1:
        text = text[a:b + 1]
    return json.loads(text)


def call_claude(client, model, prompt, max_tokens=8000):
    kwargs = dict(model=model, max_tokens=max_tokens,
                  messages=[{"role": "user", "content": prompt}])
    try:
        # temperature=0 gives determinism on models that still support it
        msg = client.messages.create(temperature=0, **kwargs)
    except Exception as e:
        if "temperature" in str(e).lower():
            msg = client.messages.create(**kwargs)  # newer models: omit it
        else:
            raise
    return "".join(b.text for b in msg.content if getattr(b, "type", "") == "text")


# ----------------------------------------------------------------- prompts
# Templates use plain @@TOKEN@@ placeholders filled via str.replace (NOT % or
# .format), so literal % and { } in the prompts can never break formatting.
EXTRACT_PROMPT = """You are an ESG analyst doing the EXTRACT step. Pull ONLY facts that are actually written in the due-diligence (DD) report below. The text is split by slide ("--- Slide N ---"); cite the slide number as the source for every fact.

CRITICAL RULE: never infer, estimate or invent anything. If a data point is NOT in the report, write the EXACT string "@@TOCONFIRM@@" as its value and add a short description of it to "data_gaps". Numbers especially must come straight from the report.

Return ONLY valid JSON, EXACTLY this shape:
{
 "company_profile": {"name":"...","sector":"...","business_model":"...","locations":"...","revenue":"...","employees":"..."},
 "company_profile_sources": {"name":"Slide N","sector":"Slide N","business_model":"Slide N","locations":"Slide N","revenue":"Slide N","employees":"Slide N"},
 "overall": {"abstain_from_deal":"Yes or No","headline":"one-sentence verdict from the report","source":"Slide N"},
 "existing_policies": [{"name":"policy/process name","status":"in place / planned / absent","source":"Slide N"}],
 "key_metrics": [{"name":"metric","value":"as stated","benchmark":"vs peer/sector if given else @@TOCONFIRM@@","source":"Slide N"}],
 "themes": [{"code":"E1","name":"theme name","pillar":"E/S/G","maturity":"rating word from report or @@TOCONFIRM@@",
             "finding":"2-3 sentences of the assessed risk/gap","recommended_actions":["action as stated in DD"],
             "investment_eur": <theme short-term investment in EUR or null>,
             "value_creation_eur": <theme annual value-creation opp in EUR or null>,"source":"Slide N"}],
 "totals": {"investment_eur": <report's STATED total investment EUR or null>,
            "value_creation_eur": <report's STATED total value-creation EUR or null>, "source":"Slide N"},
 "data_gaps": ["each notable fact NOT found in the DD"]
}
Convert 'k'->thousands, 'm'/'M'->millions. Include every material theme.

DD REPORT:
@@REPORT@@
"""

ADVISE_PROMPT = """You are a private equity value-creation expert doing the ADVISE step. Using ONLY the extracted findings below as the factual base, design an EXTENSIVE, DETAILED 100-day post-acquisition ESG plan. This is your recommendation layer - go well beyond restating the DD: add concrete activities, owners, milestones and recommended targets that a portfolio operating team could execute.

RULES:
- Everything here is a RECOMMENDATION built on the facts; do not contradict the facts.
- For any KPI baseline that is not in the findings, write the EXACT string "@@TOCONFIRM@@" - never invent a baseline number. Targets are recommendations and may be expressed as relative goals (e.g. "20 percent below baseline").
- Reference the source finding for each initiative via its theme code / slide.
- Create MULTIPLE initiatives per theme where useful, and give each 2-4 concrete activities.

Assign each initiative to ONE phase: "0-30" (stabilise & assign ownership), "30-60" (quick wins), "60-100" (embed KPIs & board reporting). Use dependency logic.
Score each 1-5 on: risk_reduction, ebitda_impact, revenue_growth, financing_benefit, operational_efficiency, implementation_feasibility.
Owner roles: "CFO","Head of Operations","Head of HR","Head of IT / Security","ESG Lead","General Counsel". Do NOT include euro figures (handled separately).

Return ONLY valid JSON, EXACTLY:
{
 "value_story": "3-4 sentence commercial narrative: how this plan protects and creates value (recommendation)",
 "assumptions": ["key assumptions behind the recommendations"],
 "initiatives": [
   {"id":"E1-01","theme":"E1","theme_name":"...","pillar":"E","title":"imperative title",
    "objective":"what this achieves and why it matters commercially",
    "activities":["concrete step 1","step 2","step 3"],
    "milestone":"what 'done' looks like by the end of the phase",
    "owner_role":"ESG Lead",
    "recommended_kpi":{"name":"...","baseline":"value from findings or @@TOCONFIRM@@","target":"recommended target"},
    "scores":{"risk_reduction":4,"ebitda_impact":3,"revenue_growth":1,"financing_benefit":2,"operational_efficiency":4,"implementation_feasibility":5},
    "phase":"0-30","dependencies":[],"source":"E1 / Slide N","rationale":"one sentence"}
 ]
}

EXTRACTED FINDINGS:
@@FINDINGS@@
"""

BOARD_PROMPT = """Prepare a concise, commercially framed board update on the 100-day ESG plan below. Return ONLY valid JSON:
{"executive_summary":"3-4 sentences","phases":[{"phase":"0-30","status":"On track/At risk/Complete","highlights":["..."]},{"phase":"30-60","status":"...","highlights":["..."]},{"phase":"60-100","status":"...","highlights":["..."]}],"top_priorities":["..."],"risks_and_asks":["..."]}
PLAN + TRACKING + METRICS:
@@PLAN@@
"""


def extract_prompt(report_text):
    return EXTRACT_PROMPT.replace("@@TOCONFIRM@@", TOCONFIRM).replace("@@REPORT@@", report_text)


def advise_prompt(findings_json):
    return ADVISE_PROMPT.replace("@@TOCONFIRM@@", TOCONFIRM).replace("@@FINDINGS@@", findings_json)


def board_prompt(plan_json):
    return BOARD_PROMPT.replace("@@PLAN@@", plan_json)


# ------------------------------------------------------------- scoring / math
def _num(x):
    try:
        return float(x)
    except (TypeError, ValueError):
        return 0.0


def compute_composite(i, w):
    s = i.get("scores", {}); tw = sum(w.values()) or 1
    return round(sum(_num(s.get(k)) * w[k] for k in w) / tw, 2)


def attach_theme_money(extract, initiatives):
    money = {t.get("code"): (t.get("investment_eur"), t.get("value_creation_eur"))
             for t in extract.get("themes", [])}
    by = {}
    for i in sorted(initiatives, key=lambda x: x.get("id", "")):
        by.setdefault(i.get("theme"), []).append(i)
    for theme, items in by.items():
        inv, val = money.get(theme, (None, None))
        for idx, it in enumerate(items):
            it["investment_eur"] = inv if idx == 0 else None
            it["value_creation_eur"] = val if idx == 0 else None
    return initiatives


def compute_metrics(extract, initiatives, tracking):
    totals = extract.get("totals") or {}; themes = extract.get("themes", [])
    tv = _num(totals.get("value_creation_eur")) or sum(_num(t.get("value_creation_eur")) for t in themes)
    ti = _num(totals.get("investment_eur")) or sum(_num(t.get("investment_eur")) for t in themes)
    roi = (tv / ti) if ti else 0.0
    done = sum(1 for i in initiatives if tracking.get(i["id"])); n = len(initiatives)
    return {"total_value_creation_eur": tv, "total_investment_eur": ti, "roi_multiple": round(roi, 1),
            "initiatives_total": n, "initiatives_complete": done,
            "percent_complete": round(100 * done / n) if n else 0}


def build_plan_object(extract, advice, weights, tracking):
    inits = advice.get("initiatives", [])
    attach_theme_money(extract, inits)
    ranked = sorted(inits, key=lambda i: compute_composite(i, weights), reverse=True)
    for r, i in enumerate(ranked, 1):
        i["composite_score"] = compute_composite(i, weights); i["priority_rank"] = r
        i["status"] = "complete" if tracking.get(i["id"]) else "open"
    cp = extract.get("company_profile", {})
    return {
        "company": {"name": cp.get("name", "Company")},
        "themes": extract.get("themes", []),
        "company_profile": cp, "company_profile_sources": extract.get("company_profile_sources", {}),
        "overall": extract.get("overall", {}), "existing_policies": extract.get("existing_policies", []),
        "key_metrics": extract.get("key_metrics", []), "data_gaps": extract.get("data_gaps", []),
        "value_story": advice.get("value_story", ""), "assumptions": advice.get("assumptions", []),
        "generated_on": str(date.today()), "weights": weights,
        "metrics": compute_metrics(extract, inits, tracking),
        "phases": {k: {"label": l, "initiative_ids": [i["id"] for i in ranked if i.get("phase") == k]}
                   for k, l in PHASES},
        "initiatives": ranked,
    }


# --------------------------------------------------------- pptx helpers
def _set(run, size, bold, color, italic=False):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color


def _box(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    return tf


def _first(tf, t, s, bold=False, color=TEXT, italic=False):
    tf.text = t; _set(tf.paragraphs[0].runs[0], s, bold, color, italic); return tf.paragraphs[0]


def _line(tf, t, s, bold=False, color=TEXT, italic=False, before=0, lvl=0):
    p = tf.add_paragraph(); p.text = t; p.space_before = Pt(before); p.level = lvl
    col = REDFLAG if (isinstance(t, str) and "TO CONFIRM" in t) else color
    _set(p.runs[0], s, bold, col, italic); return p


def _rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); return sh


def _logo(slide):
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(LOGO_PATH, Inches(11.0), Inches(0.35), width=Inches(1.9))
        except Exception:
            pass


def _footer(slide, note=""):
    tf = _box(slide, 0.5, 7.12, 12.3, 0.3)
    txt = "Strictly Private & Confidential. Copyright \u00a9 %d Holtara" % date.today().year
    if note:
        txt = note + "   |   " + txt
    _first(tf, txt, 8, color=GREY)


def _content(prs, title, sub=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0.0, 0.0, 13.33, 0.95, NAVY)
    h = _box(s, 0.5, 0.12, 10.3, 0.8)
    _first(h, title, 22, bold=True, color=WHITE)
    if sub:
        _line(h, sub, 11, color=RGBColor(0xCF, 0xD6, 0xE6))
    _logo(s)
    return s


# ----------------------------------------------------------- DETAILED deck
def detailed_plan_pptx(plan):
    prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    eur = lambda v: "EUR {:,.0f}".format(_num(v))
    cp = plan.get("company_profile", {}); cps = plan.get("company_profile_sources", {})
    company = plan.get("company", {}).get("name", "Company"); m = plan.get("metrics", {})

    # 1 Title
    s = prs.slides.add_slide(blank)
    _rect(s, 0, 2.4, 13.33, 2.7, NAVY)
    t = _box(s, 0.7, 2.7, 11, 2.0)
    _first(t, "100-Day ESG Value Creation Plan", 38, bold=True, color=WHITE)
    _line(t, "Driving positive change", 14, color=RGBColor(0xCF, 0xD6, 0xE6), italic=True)
    _line(t, company, 20, color=ORANGE)
    _logo(s); _footer(s)

    # 2 Executive summary
    s = _content(prs, "Executive summary", "The value story (recommendation) and headline economics")
    cards = [("Annual value creation", eur(m.get("total_value_creation_eur", 0)), ORANGE),
             ("Investment", eur(m.get("total_investment_eur", 0)), NAVY),
             ("Return on investment", "%sx" % m.get("roi_multiple", 0), ORANGE),
             ("Initiatives", str(m.get("initiatives_total", 0)), NAVY)]
    x = 0.5
    for label, val, col in cards:
        c = _box(s, x, 1.2, 3.05, 1.0); _first(c, val, 23, bold=True, color=col)
        _line(c, label, 10.5, color=GREY); x += 3.07
    vs = _box(s, 0.5, 2.45, 12.3, 1.7)
    _first(vs, "Value story", 13, bold=True, color=NAVY)
    _line(vs, plan.get("value_story", ""), 12, color=TEXT)
    ov = plan.get("overall", {})
    vd = _box(s, 0.5, 4.2, 12.3, 1.6)
    _first(vd, "Overall ESG verdict (from DD)", 13, bold=True, color=NAVY)
    _line(vd, "%s   (source: %s)" % (ov.get("headline", TOCONFIRM), ov.get("source", "n/a")), 12, color=TEXT)
    _footer(s)

    # 3 Baseline facts (EXTRACT, traceable)
    s = _content(prs, "Company & ESG baseline", "Facts extracted from the DD \u2014 source shown in brackets")
    left = _box(s, 0.5, 1.2, 6.2, 5.7)
    _first(left, "Company profile", 13, bold=True, color=ORANGE)
    for k in ["sector", "business_model", "locations", "revenue", "employees"]:
        _line(left, "%s: %s" % (k.replace("_", " ").title(), cp.get(k, TOCONFIRM)), 11, before=3)
        _line(left, "source: %s" % cps.get(k, "n/a"), 8.5, color=GREY)
    _line(left, "Existing policies & processes", 13, bold=True, color=ORANGE, before=10)
    for p in plan.get("existing_policies", [])[:6]:
        _line(left, "\u2022 %s \u2014 %s  (%s)" % (p.get("name", ""), p.get("status", ""), p.get("source", "")), 10.5, before=2)
    right = _box(s, 7.0, 1.2, 5.8, 5.7)
    _first(right, "Key metrics (as reported)", 13, bold=True, color=ORANGE)
    for met in plan.get("key_metrics", [])[:9]:
        _line(right, "\u2022 %s: %s" % (met.get("name", ""), met.get("value", TOCONFIRM)), 10.5, before=3)
        bm = met.get("benchmark", "")
        _line(right, "benchmark: %s  (%s)" % (bm, met.get("source", "")), 8.5, color=GREY)
    _footer(s, "Traceable to DD")

    # 4-6 Phase detail
    for pkey, plabel in PHASES:
        items = [i for i in plan.get("initiatives", []) if i.get("phase") == pkey]
        s = _content(prs, plabel.replace("\u00b7", "\u2014"), "%d initiative(s) \u00b7 detailed actions, owners, milestones & KPIs" % len(items))
        body = _box(s, 0.5, 1.15, 12.3, 5.9)
        if not items:
            _first(body, "No initiatives in this phase.", 12, color=GREY)
        else:
            shown = items[:4]
            firstdone = False
            for it in shown:
                title = "%s   \u00b7   %s   \u00b7   priority %.1f" % (it.get("title", ""), it.get("owner_role", "-"), it.get("composite_score", 0))
                if not firstdone:
                    _first(body, title, 13, bold=True, color=PHASE_COLOR[pkey]); firstdone = True
                else:
                    _line(body, title, 13, bold=True, color=PHASE_COLOR[pkey], before=10)
                _line(body, it.get("objective", ""), 10.5, color=TEXT)
                for a in it.get("activities", [])[:4]:
                    _line(body, "\u2022 " + a, 10, color=TEXT, before=1, lvl=1)
                kpi = it.get("recommended_kpi", {}) or {}
                _line(body, "KPI: %s  |  baseline: %s  \u2192  target: %s" % (
                    kpi.get("name", "-"), kpi.get("baseline", TOCONFIRM), kpi.get("target", "-")), 9.5, color=NAVY, before=2)
                extra = []
                if it.get("milestone"):
                    extra.append("Milestone: " + it["milestone"])
                if it.get("value_creation_eur"):
                    extra.append("Value: +%s/yr" % eur(it.get("value_creation_eur")))
                extra.append("source: %s" % it.get("source", "n/a"))
                _line(body, "   ".join(extra), 8.5, color=GREY)
            if len(items) > 4:
                _line(body, "+ %d more initiative(s) in this phase (see JSON export)" % (len(items) - 4), 9.5, color=GREY, before=8)
        _footer(s)

    # 7 Targets & value story (ADVISE)
    s = _content(prs, "Recommended targets & value story", "Holtara recommendations built on the DD findings")
    body = _box(s, 0.5, 1.2, 12.3, 5.7)
    _first(body, "Recommended targets by initiative", 13, bold=True, color=ORANGE)
    for it in plan.get("initiatives", []):
        kpi = it.get("recommended_kpi", {}) or {}
        if kpi.get("target"):
            _line(body, "\u2022 %s: %s \u2192 %s" % (it.get("title", ""), kpi.get("baseline", TOCONFIRM), kpi.get("target", "")), 10.5, before=2)
    _line(body, "Assumptions behind these recommendations", 13, bold=True, color=ORANGE, before=12)
    for a in plan.get("assumptions", []):
        _line(body, "\u2022 " + a, 10.5, before=2)
    _footer(s, "Recommendations")

    # 8 Data gaps
    s = _content(prs, "Data to confirm", "Not found in the DD \u2014 confirm before finalising (no numbers were invented)")
    body = _box(s, 0.5, 1.2, 12.3, 5.7)
    gaps = plan.get("data_gaps", [])
    if not gaps:
        _first(body, "No material data gaps identified.", 12, color=GREY)
    else:
        _first(body, "%d item(s) flagged as %s" % (len(gaps), TOCONFIRM), 12, bold=True, color=REDFLAG)
        for g in gaps:
            _line(body, "\u2022 " + g, 11, before=3)
    _footer(s)

    out = BytesIO(); prs.save(out); return out.getvalue()


# ----------------------------------------------------------- one-page summary
def onepager_pptx(plan):
    prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    eur = lambda v: "EUR {:,.0f}".format(_num(v))
    company = plan.get("company", {}).get("name", "Company"); m = plan.get("metrics", {})
    head = _box(s, 0.5, 0.35, 9.5, 1.1)
    _first(head, "100-Day ESG Value Creation Plan", 26, bold=True, color=ORANGE)
    _line(head, "Driving positive change", 12, color=TEAL, italic=True)
    _line(head, company, 13, color=NAVY)
    _logo(s); _rect(s, 0.5, 1.55, 12.33, 0.03, ORANGE)
    cards = [("Annual value creation", eur(m.get("total_value_creation_eur", 0)), ORANGE),
             ("Investment", eur(m.get("total_investment_eur", 0)), NAVY),
             ("Return on investment", "%sx" % m.get("roi_multiple", 0), ORANGE),
             ("Initiatives", str(m.get("initiatives_total", 0)), NAVY)]
    x = 0.5
    for label, val, col in cards:
        c = _box(s, x, 1.75, 3.05, 1.0); _first(c, val, 24, bold=True, color=col)
        _line(c, label, 10.5, color=GREY); x += 3.07
    _rect(s, 0.5, 2.85, 12.33, 0.02, RGBColor(0xDD, 0xDD, 0xDB))
    for (pkey, plabel), cx in zip(PHASES, [0.5, 4.78, 9.06]):
        _rect(s, cx, 3.02, 4.05, 0.42, PHASE_COLOR[pkey])
        hb = _box(s, cx + 0.12, 3.05, 3.85, 0.4)
        _first(hb, plabel.split("\u00b7")[0].strip(), 12, bold=True, color=WHITE)
        items = [i for i in plan.get("initiatives", []) if i.get("phase") == pkey]
        body = _box(s, cx + 0.05, 3.55, 4.0, 3.35)
        _first(body, (plabel.split("\u00b7")[1].strip() if "\u00b7" in plabel else "").upper(), 8.5, bold=True, color=GREY)
        for it in items[:6]:
            money = "  \u00b7  +%s/yr" % eur(it.get("value_creation_eur")) if it.get("value_creation_eur") else ""
            _line(body, "\u2022 " + it.get("title", ""), 10.5, bold=True, color=NAVY, before=6)
            _line(body, "%s%s" % (it.get("owner_role", ""), money), 8.5, color=GREY)
        if len(items) > 6:
            _line(body, "+ %d more" % (len(items) - 6), 9, color=GREY, before=4)
    _footer(s)
    out = BytesIO(); prs.save(out); return out.getvalue()


# ----------------------------------------------------------- FINDINGS one-pager
def findings_onepager_pptx(plan):
    """Single branded slide summarising the EXTRACTED DD findings."""
    prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    eur = lambda v: "EUR {:,.0f}".format(_num(v))
    company = plan.get("company", {}).get("name", "Company")
    m = plan.get("metrics", {}); ov = plan.get("overall", {})
    themes = plan.get("themes", []); gaps = plan.get("data_gaps", [])

    head = _box(s, 0.5, 0.35, 9.5, 1.1)
    _first(head, "ESG Due Diligence \u2014 Findings", 26, bold=True, color=ORANGE)
    _line(head, "Driving positive change", 12, color=TEAL, italic=True)
    _line(head, company, 13, color=NAVY)
    _logo(s); _rect(s, 0.5, 1.55, 12.33, 0.03, ORANGE)

    # verdict + headline cards (from the DD)
    abstain = ov.get("abstain_from_deal", TOCONFIRM)
    cards = [("Value-creation opportunity", eur(m.get("total_value_creation_eur", 0)), ORANGE),
             ("Investment indicated", eur(m.get("total_investment_eur", 0)), NAVY),
             ("Material themes", str(len(themes)), ORANGE),
             ("Abstain from deal?", abstain, REDFLAG if str(abstain).lower().startswith("y") else NAVY)]
    x = 0.5
    for label, val, col in cards:
        c = _box(s, x, 1.72, 3.05, 0.95); _first(c, val, 22, bold=True, color=col)
        _line(c, label, 10, color=GREY); x += 3.07

    vd = _box(s, 0.5, 2.7, 12.3, 0.5)
    _first(vd, "Overall verdict: %s   (source: %s)" % (ov.get("headline", TOCONFIRM), ov.get("source", "n/a")), 11, color=TEXT, italic=True)
    _rect(s, 0.5, 3.18, 12.33, 0.02, RGBColor(0xDD, 0xDD, 0xDB))

    # findings grid (2 columns)
    _box(s, 0.5, 3.24, 6, 0.3) and _first(_box(s, 0.5, 3.24, 8, 0.3), "Material ESG findings (extracted from the DD)", 11, bold=True, color=GREY)
    cols_x = [0.5, 6.85]
    for idx, t in enumerate(themes[:6]):
        col, row = idx % 2, idx // 2
        cx = cols_x[col]; cy = 3.6 + row * 1.18
        pc = PILLAR_COLOR.get(t.get("pillar"), NAVY)
        _rect(s, cx, cy, 0.08, 1.02, pc)
        tb = _box(s, cx + 0.18, cy - 0.02, 5.75, 1.12)
        _first(tb, "%s \u00b7 %s   (%s)" % (t.get("code", ""), t.get("name", ""), t.get("maturity", TOCONFIRM)), 11, bold=True, color=pc)
        fnd = t.get("finding", "")
        _line(tb, fnd if len(fnd) < 165 else fnd[:162] + "...", 9.5, color=TEXT)
        _line(tb, "source: %s" % t.get("source", "n/a"), 8, color=GREY)
    if len(themes) > 6:
        _first(_box(s, 0.5, 6.75, 6, 0.3), "+ %d more theme(s) in the detailed plan" % (len(themes) - 6), 9, color=GREY)

    note = "Data to confirm: %d item(s) not in the DD \u2014 see detailed plan." % len(gaps) if gaps else ""
    _footer(s, note)
    out = BytesIO(); prs.save(out); return out.getvalue()


# ----------------------------------------------------------- board update
def board_markdown(company_name, n, m):
    eur = lambda v: "EUR {:,.0f}".format(_num(v))
    L = ["# 100-Day ESG Value Creation Plan - Board Update",
         "**Company:** %s  |  **Date:** %s" % (company_name, date.today()), "",
         "## Executive summary", n.get("executive_summary", ""), "", "## Key metrics",
         "- Annual value creation: **%s**" % eur(m["total_value_creation_eur"]),
         "- Investment: **%s**" % eur(m["total_investment_eur"]),
         "- ROI: **%sx**" % m["roi_multiple"],
         "- Initiatives: **%d total, %d complete (%d%%)**" % (m["initiatives_total"], m["initiatives_complete"], m["percent_complete"]),
         "", "## Progress by phase"]
    for ph in n.get("phases", []):
        L.append("**%s - %s**" % (ph.get("phase", ""), ph.get("status", "")))
        L += ["- %s" % h for h in ph.get("highlights", [])]; L.append("")
    L.append("## Top priorities"); L += ["- %s" % p for p in n.get("top_priorities", [])]
    L.append(""); L.append("## Risks & asks"); L += ["- %s" % r for r in n.get("risks_and_asks", [])]
    return "\n".join(L)


def board_pptx(company_name, n, m):
    prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    s = _content(prs, "Board update \u2014 100-Day ESG Plan", "%s  \u00b7  %s" % (company_name, date.today()))
    eur = lambda v: "EUR {:,.0f}".format(_num(v))
    cards = [("Value / yr", eur(m["total_value_creation_eur"]), ORANGE), ("Investment", eur(m["total_investment_eur"]), NAVY),
             ("ROI", "%sx" % m["roi_multiple"], ORANGE), ("Progress", "%d%%" % m["percent_complete"], NAVY)]
    x = 0.5
    for label, val, col in cards:
        c = _box(s, x, 1.2, 3.05, 1.0); _first(c, val, 22, bold=True, color=col); _line(c, label, 10.5, color=GREY); x += 3.07
    body = _box(s, 0.5, 2.4, 12.3, 4.4)
    _first(body, "Executive summary", 13, bold=True, color=NAVY)
    _line(body, n.get("executive_summary", ""), 12)
    for ph in n.get("phases", []):
        _line(body, "%s  -  %s" % (ph.get("phase", ""), ph.get("status", "")), 12, bold=True, color=NAVY, before=8)
        for h in ph.get("highlights", []):
            _line(body, "\u2022 " + h, 11, before=1)
    _footer(s)
    out = BytesIO(); prs.save(out); return out.getvalue()
